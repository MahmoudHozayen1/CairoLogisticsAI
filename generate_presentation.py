"""Generate the project presentation (docs/presentation.pptx) with python-pptx.

Run::

    python generate_presentation.py

The deck mirrors docs/PRESENTATION_OUTLINE.md. Edit the SLIDES list below (or the
outline) and regenerate.
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- Brand palette (matches the web app) ---
PRIMARY = RGBColor(0x0D, 0x3B, 0x66)
ACCENT = RGBColor(0xEE, 0x6C, 0x4D)
LIGHT = RGBColor(0xF4, 0xF7, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x29, 0x33)

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "docs", "presentation.pptx")
BENCH_DIR = os.path.join(HERE, "docs", "benchmarks")

# (title, [bullets], subtitle_or_None)
SLIDES = [
    ("SwiftRoute", [
        "AI-Assisted Last-Mile Logistics & Parcel Delivery Platform",
        "A full-stack system inspired by Bosta & DHL",
        "Merchants ship \u2192 AI optimises routes \u2192 couriers deliver \u2192 customers track",
        "Every automated decision explains its reasoning",
    ], "TITLE"),
    ("The Problem", [
        "Last-mile delivery is the most expensive, least-optimised leg of the supply chain.",
        "Cairo: dense congestion, landmark-based addresses, cash-on-delivery everywhere.",
        "Small merchants coordinate couriers manually by phone \u2014 no routing, no ETAs, no proof.",
        "The OR/ML techniques that could fix this are locked in proprietary black boxes.",
    ], None),
    ("What We Built (at a glance)", [
        "Multi-role platform \u2014 Admin, Courier, Merchant portals + public no-login tracking.",
        "AI route optimiser \u2014 k-means clustering + 6 TSP techniques, benchmarked.",
        "Explainable ML suite \u2014 5 slices: ETA, forecasting, neural router, personas, NLP + assistant.",
        "Trust layer \u2014 SHA-256 hash-chained custody + GIS proof of delivery.",
        "Street-following routes \u2014 real OSRM geometry, simulated traffic, road closures.",
        "74 automated tests; deployable one-click on a free tier.",
    ], None),
    ("The Big-Picture Pipeline", [
        "Merchant drops a pin on the map \u2192 shipment created.",
        "Parcel lands at its hub (at_warehouse).",
        "AI optimiser: k-means clusters \u2192 TSP sequence \u2192 OSRM street geometry + ETA.",
        "Courier delivers with photo + GPS \u2192 hash-chain custody + GIS proof.",
        "Customer follows a live public tracking timeline.",
        "The optimiser and every model expose their reasoning throughout.",
    ], None),
    ("Architecture", [
        "Flask application factory + blueprints (main, auth, admin, courier, merchant, tracking, api).",
        "SQLAlchemy ORM \u2014 SQLite zero-config by default, PostgreSQL via one DATABASE_URL.",
        "routing/ = optimizer + street_router; ml/ = predictive layer; audit.py = custody + GIS.",
        "Front-end: Jinja2 + Bootstrap 5 + Leaflet maps + Chart.js dashboards.",
        "Security: PBKDF2 hashing, CSRF on every form, role_required decorator, open-redirect guard.",
    ], None),
    ("The Optimisation Engine \u2014 4 Steps", [
        "Vehicle routing is NP-hard \u2014 we decompose it into tractable sub-problems.",
        "1) Assign \u2014 each parcel belongs to its hub.",
        "2) Cluster \u2014 k-means gives each courier a zone; capacity-aware (Bike 3 / Moto 5 / Car 10 / Van 15).",
        "3) Sequence \u2014 solve the TSP order with a selectable technique.",
        "4) Geometry + ETA \u2014 OSRM real-street paths, cached, straight-line fallback.",
    ], None),
    ("Six TSP Techniques", [
        "FIFO \u2014 as received: the baseline to beat.",
        "Nearest Neighbour \u2014 fast greedy heuristic.",
        "2-opt \u2014 un-cross edges; the default (best quality/speed balance).",
        "Or-opt \u2014 adds segment relocation; highest quality.",
        "Christofides (NetworkX) \u2014 proven 1.5\u00d7-optimal guarantee for metric TSP.",
        "Simulated Annealing (NetworkX) \u2014 metaheuristic escaping local optima.",
    ], None),
    ("\u2b50 WOW: Benchmarked, Not Asserted", [
        "Reproducible harness; brute-force exact optimum for size \u2264 8.",
        "Optimised routes cut total distance ~60\u201362% vs naive FIFO.",
        "Or-opt lands ~0.1% from the exact optimum \u2014 in milliseconds; 2-opt & SA < 1%.",
        "Design defended by data, not opinion.",
    ], "improvement_by_strategy.png"),
    ("\u2b50 WOW: Optimality & Speed", [
        "On solvable instances 2-opt & SA stay < 1% above optimal; Christofides ~9% but guaranteed.",
        "2-opt gives the best quality/speed balance; runtimes scale gently.",
    ], "runtime_scaling.png"),
    ("\u2b50 WOW: Live Traffic & Road Closures", [
        "Routes coloured green \u2192 amber \u2192 orange \u2192 red by congestion.",
        "Deterministic time-of-day + day-of-week model (rush-hour curve, Friday lightest for Cairo).",
        "Honest engineering: no fake live feed \u2014 clearly labelled simulated in the UI.",
        "Closures are first-class: the optimiser re-routes around them; un-avoidable legs flagged dashed red.",
        "Smart split: store expensive geometry, compute cheap traffic fresh at render time.",
    ], None),
    ("\u2b50 WOW: Dispatch-Time Planning", [
        "\u201cIs Friday morning faster than Tuesday 6pm?\u201d \u2014 now answerable.",
        "traffic_factor_at() scales each leg's ETA by the traffic when the courier actually reaches it.",
        "compare_strategies() ranks techniques by makespan (when the last parcel lands), then distance.",
        "Preview is pure estimation \u2014 instant, offline, no DB writes; only dispatch persists.",
    ], None),
    ("The ML Suite \u2014 5 Slices", [
        "All pure numpy / scikit-learn \u2014 no torch, no GPU, trains in seconds; deploys on a free tier.",
        "1) ETA + forecasting  2) feedback + audit + GIS  3) NLP notes + assistant.",
        "4) learning-to-route neural policy  5) courier-behaviour personas.",
        "Trained on a seeded, reproducible synthetic history (~5.8k deliveries / 180 days).",
        "Guiding principle: every model explains itself.",
    ], None),
    ("Slice 1 \u2014 ETA & Demand Forecasting", [
        "Gradient-boosted trees: drop-off ETA (MAE ~2.7 min, R\u00b2 ~0.94), pickup ETA (MAE ~1.3 min).",
        "Late-risk classifier: AUC ~0.82 (5-fold CV, 72% better than mean baseline).",
        "SeasonalTrendForecaster: trend + weekly seasonality + 95% band, MAPE ~14%.",
        "TreeContributionExplainer gives exact additive \u2018why\u2019 per prediction \u2014 no black box.",
    ], None),
    ("Slice 2 \u2014 Trust & Audit", [
        "Chain of custody: every handoff SHA-256 hash-linked to the previous \u2014 tamper-evident.",
        "verify_chain() recomputes and detects tampering.",
        "GIS proof: haversine geofence confirms delivery within radius of the destination.",
        "Feedback/drift monitor: predicted-vs-actual log, weekly drift flags, calibration + Brier score.",
    ], None),
    ("Slice 3 \u2014 NLP Notes + Assistant", [
        "Hybrid analyser: lexicon/regex + TF-IDF one-vs-rest logistic regression (10-tag taxonomy).",
        "Generalises to unseen phrasing (\u201cthe vase can shatter\u201d \u2192 fragile), with exact token reasoning.",
        "Extracts time windows (\u201cbetween 6 and 9 pm\u201d \u2192 18:00\u201321:00) and target floor.",
        "Assistant: DB retrieval always works + optional local LLM that only phrases the facts.",
    ], None),
    ("Slice 4 \u2014 Learning-to-Route (Neural)", [
        "Pure-numpy pointer-network policy; linear-softmax attention over unvisited stops (7 features).",
        "Trained with REINFORCE + a greedy-rollout baseline (self-critical / Kool-style).",
        "Key insight: greedy collapses to nearest-neighbour \u2014 the win is best-of-N sampled rollouts.",
        "Result: ~5% shorter than NN, ~1.5\u20132% gap to 2-opt; per-step feature reasoning.",
    ], None),
    ("Slice 5 \u2014 Courier-Behaviour Personas", [
        "Simulated GPS traces \u2192 stop detection (idle / delivery / break by dwell time).",
        "7 engineered features \u2192 k-means into 4 personas: Efficient / Steady / Idle-prone / At-risk.",
        "Validated: silhouette 0.36, ARI 0.99 vs ground-truth archetypes.",
        "z-score reasoning + productivity score \u2014 fleet insight with no new hardware.",
    ], None),
    ("Hardest Difficulties We Solved", [
        "Behaviour bug: dwell measured from arrival ping, not within-run span \u2014 fixed all-idle misclassification.",
        "Hash chain vs SQLite: timestamps must be naive UTC or the chain won't re-verify after round-trip.",
        "Explainer precision: walk trees in float32 to match sklearn's leaf selection exactly.",
        "Config trap: a global ~/.env DATABASE_URL hijacked the app \u2192 made config project-local.",
        "Deploy bug: AUTO_INIT_DB=true parsed as false \u2192 no tables \u2192 500 on login.",
    ], None),
    ("Key Engineering Decisions", [
        "Explainability first \u2014 every model returns signed, human-readable reasoning.",
        "Deploy-anywhere \u2014 light deps only; every heavy library/network call has a fallback.",
        "Additive schema \u2014 new features = new tables or simulation; create_all upgrades old DBs cleanly.",
        "Reproducible \u2014 everything seeded (seed=42).",
        "Key-less philosophy (OSRM + Nominatim); map-pin addressing beats landmark geocoding.",
    ], None),
    ("Quality & Validation", [
        "74 automated tests \u2014 lifecycle, optimiser, ML, audit, NLP, router, behaviour.",
        "Tests run offline (ROUTING_PROVIDER=straight).",
        "Models validated against baselines with cross-validation \u2014 honest, not lucky-split.",
        "Reproducible benchmark harness with exact-optimum ground truth.",
    ], None),
    ("Impact & Future Work", [
        "Impact: an open, explainable alternative to proprietary logistics black boxes on minimal infra.",
        "Real operational data + a live traffic feed.",
        "Hard time windows (VRPTW) + multi-day tours; live mid-trip re-routing from real GPS.",
        "Full attention/transformer router on capable hardware; mobile app + OTP-verified handoffs.",
    ], None),
    ("Thank You / Demo", [
        "Try it: Merchant creates a shipment \u2192 Admin optimises \u2192 Courier delivers \u2192 track publicly.",
        "Demo: admin@swiftroute.app / courier1@swiftroute.app / merchant1@swiftroute.app",
        "Questions & discussion.",
    ], "TITLE"),
]


def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _title_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, PRIMARY)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = title
    run.font.size = Pt(54); run.font.bold = True; run.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    for b in bullets:
        para = tf.add_paragraph()
        r = para.add_run(); r.text = b
        r.font.size = Pt(20); r.font.color.rgb = RGBColor(0xCF, 0xDD, 0xEC)
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(10)
    # accent bar
    bar = slide.shapes.add_shape(1, Inches(5.4), Inches(1.7), Inches(2.5), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()


def _content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, LIGHT)
    # header band
    band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.25))
    band.fill.solid(); band.fill.fore_color.rgb = PRIMARY; band.line.fill.background()
    tbox = band.text_frame; tbox.word_wrap = True
    tbox.margin_left = Inches(0.6)
    p = tbox.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.4))
    tf = body.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet = para.add_run(); bullet.text = "\u25B8  "
        bullet.font.size = Pt(20); bullet.font.color.rgb = ACCENT; bullet.font.bold = True
        run = para.add_run(); run.text = b
        run.font.size = Pt(20); run.font.color.rgb = DARK
        para.space_after = Pt(14)


def _image_slide(prs, title, bullets, image_filename):
    """A content slide whose lower half showcases a benchmark chart image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, LIGHT)
    band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.25))
    band.fill.solid(); band.fill.fore_color.rgb = PRIMARY; band.line.fill.background()
    tbox = band.text_frame; tbox.word_wrap = True
    tbox.margin_left = Inches(0.6)
    p = tbox.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.45), Inches(11.5), Inches(1.4))
    tf = body.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet = para.add_run(); bullet.text = "\u25B8  "
        bullet.font.size = Pt(16); bullet.font.color.rgb = ACCENT; bullet.font.bold = True
        run = para.add_run(); run.text = b
        run.font.size = Pt(16); run.font.color.rgb = DARK
        para.space_after = Pt(6)

    path = os.path.join(BENCH_DIR, image_filename)
    if os.path.exists(path):
        pic_w = Inches(7.47)  # keeps the 8:4.5 chart aspect ratio at 4.2" tall
        left = (Inches(13.333) - pic_w) / 2
        slide.shapes.add_picture(path, left, Inches(2.9), height=Inches(4.2))


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for title, bullets, kind in SLIDES:
        if kind == "TITLE":
            _title_slide(prs, title, bullets)
        elif isinstance(kind, str) and kind.endswith(".png"):
            _image_slide(prs, title, bullets, kind)
        else:
            _content_slide(prs, title, bullets)
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"Presentation written to {OUT} ({len(SLIDES)} slides).")


if __name__ == "__main__":
    build()
